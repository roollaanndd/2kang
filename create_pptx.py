#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# Brand Colors
MAGENTA     = RGBColor(0xC9, 0x1F, 0x7A)
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)
TEAL        = RGBColor(0x00, 0xC5, 0xCC)
OFF_WHITE   = RGBColor(0xF6, 0xF7, 0xFB)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY  = RGBColor(0xE5, 0xE7, 0xEB)
LIGHT_MAG   = RGBColor(0xFC, 0xE7, 0xF3)
NAVY2       = RGBColor(0x15, 0x28, 0x55)
TEAL_DARK   = RGBColor(0x00, 0x96, 0x9B)

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill, border=None, bw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw or 0.75)
    else: s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill, border=None, bw=None):
    s = slide.shapes.add_shape(5, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw or 0.75)
    else: s.line.fill.background()
    return s

def oval(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, sz=12, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def header(slide, title, sub=None):
    rect(slide, 0, 0, W, Inches(1.45), NAVY)
    rect(slide, 0, Inches(1.42), W, Pt(4), MAGENTA)
    txt(slide, title, Inches(0.5), Inches(0.18), W-Inches(3), Inches(0.7),
        sz=26, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.5), Inches(0.88), W-Inches(3), Inches(0.5),
            sz=12, color=TEAL)
    txt(slide, "Ruang.dev × OMDC 2026", W-Inches(2.5), Inches(0.5),
        Inches(2.3), Inches(0.4), sz=9, color=RGBColor(0x9C,0xA3,0xAF),
        align=PP_ALIGN.RIGHT)
    return Inches(1.5)

def footer(slide):
    rect(slide, 0, H-Inches(0.38), W, Inches(0.38), NAVY)
    txt(slide, "OMDC Digital Transformation 2026  |  Ruang.dev  |  Confidential",
        Inches(0.5), H-Inches(0.35), W-Inches(1), Inches(0.32),
        sz=8, color=RGBColor(0x9C,0xA3,0xAF))

def bg(slide, color=OFF_WHITE):
    rect(slide, 0, 0, W, H, color)

def card(slide, l, t, w, h, fill=WHITE, border=LIGHT_GRAY, accent=MAGENTA):
    rrect(slide, l, t, w, h, fill, border, 0.5)
    rect(slide, l, t, Pt(4), h, accent)

def kpi_card(slide, l, t, w, h, number, label, sub="", nc=MAGENTA):
    rrect(slide, l, t, w, h, WHITE, LIGHT_GRAY, 0.5)
    rect(slide, l, t, w, Pt(4), nc)
    txt(slide, number, l, t+Inches(0.15), w, Inches(0.7),
        sz=40, bold=True, color=nc, align=PP_ALIGN.CENTER)
    txt(slide, label, l, t+Inches(0.85), w, Inches(0.45),
        sz=12, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, sub, l, t+Inches(1.25), w, Inches(0.35),
            sz=9, color=GRAY, align=PP_ALIGN.CENTER)

def tag(slide, text, l, t, w, h, bg_c=LIGHT_MAG, tc=MAGENTA):
    rrect(slide, l, t, w, h, bg_c)
    txt(slide, text, l, t, w, h, sz=9, bold=True, color=tc, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────
s = blank_slide()
bg(s, NAVY)
oval(s, W-Inches(5.5), Inches(-1.5), Inches(7), Inches(7), NAVY2)
oval(s, W-Inches(4), H-Inches(2.5), Inches(4), Inches(4), RGBColor(0x00,0x70,0x75))
rect(s, 0, 0, Inches(0.55), H, MAGENTA)
rect(s, 0, H-Inches(0.08), W, Inches(0.08), TEAL)

txt(s, "PRESENTED BY  RUANG.DEV", Inches(0.9), Inches(0.7),
    Inches(6), Inches(0.4), sz=10, bold=True, color=MAGENTA)
txt(s, "OMDC", Inches(0.9), Inches(1.15), Inches(9), Inches(1.3),
    sz=84, bold=True, color=WHITE)
txt(s, "DIGITAL TRANSFORMATION", Inches(0.9), Inches(2.35), Inches(10), Inches(0.8),
    sz=36, bold=True, color=MAGENTA)
txt(s, "2026", Inches(0.9), Inches(3.1), Inches(5), Inches(0.75),
    sz=36, bold=True, color=TEAL)
rect(s, Inches(0.9), Inches(3.95), Inches(8.5), Pt(2), RGBColor(0x2D,0x3F,0x7A))
txt(s, "Website Revamp  •  Mobile Ecosystem  •  Smart eKiosk Platform",
    Inches(0.9), Inches(4.1), Inches(11), Inches(0.5), sz=15, color=RGBColor(0xA0,0xB4,0xD0))
txt(s, "Disiapkan untuk OMDC Group  |  Oktri Manessa Dental Clinic",
    Inches(0.9), Inches(4.7), Inches(10), Inches(0.45), sz=12, color=GRAY)

rrect(s, W-Inches(4.2), Inches(2.0), Inches(3.2), Inches(1.7), NAVY2, MAGENTA, 1)
txt(s, "OMDC", W-Inches(4.2), Inches(2.2), Inches(3.2), Inches(0.8),
    sz=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Oktri Manessa Dental Clinic", W-Inches(4.2), Inches(3.0), Inches(3.2), Inches(0.4),
    sz=10, color=MAGENTA, align=PP_ALIGN.CENTER)

# Ecosystem icons (text-based)
eco_y = Inches(5.5)
for i, (ico, lbl) in enumerate([("🌐","Website"), ("📱","Mobile App"), ("🖥️","eKiosk")]):
    ex = Inches(0.9) + i * Inches(2.2)
    rrect(s, ex, eco_y, Inches(1.9), Inches(0.75), NAVY2, TEAL, 0.5)
    txt(s, f"{ico}  {lbl}", ex, eco_y+Inches(0.1), Inches(1.9), Inches(0.5),
        sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 2:
        txt(s, "→", ex+Inches(1.9), eco_y+Inches(0.2), Inches(0.3), Inches(0.4),
            sz=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
print("✓ Slide 1: Cover")

# ─────────────────────────────────────────
# SLIDE 2 — EXECUTIVE SUMMARY
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Ringkasan Eksekutif",
            "The Future of Dental Patient Experience — Visi Transformasi OMDC")
footer(s)

col_w = (W - Inches(1.0)) / 2
lx, rx = Inches(0.4), Inches(0.4) + col_w + Inches(0.2)
cy = ht + Inches(0.15)

# Left header
rrect(s, lx, cy, col_w-Inches(0.1), Inches(0.45), LIGHT_MAG)
txt(s, "⚠  KONDISI SAAT INI", lx+Inches(0.1), cy+Inches(0.05),
    col_w-Inches(0.2), Inches(0.35), sz=11, bold=True, color=MAGENTA)

pains = [
    "Registrasi pasien manual — 15–20 menit per pasien",
    "Antrean front desk yang panjang dan tidak terukur",
    "Tidak ada akses self-service sama sekali",
    "Data tersebar di sistem yang tidak terintegrasi",
    "Zero analytics dan laporan manajemen real-time",
]
for i, p in enumerate(pains):
    iy = cy + Inches(0.55) + i*Inches(0.6)
    rrect(s, lx, iy, col_w-Inches(0.1), Inches(0.52), WHITE, LIGHT_GRAY, 0.5)
    rect(s, lx, iy, Pt(4), Inches(0.52), MAGENTA)
    txt(s, "✗  "+p, lx+Inches(0.12), iy+Inches(0.06),
        col_w-Inches(0.3), Inches(0.42), sz=11, color=DARK_TEXT)

# Vertical divider
rect(s, lx+col_w+Inches(0.05), cy, Pt(2), H-ht-Inches(0.6), LIGHT_GRAY)

# Right header
rrect(s, rx, cy, col_w-Inches(0.1), Inches(0.45), RGBColor(0xE0,0xFB,0xFC))
txt(s, "✦  VISI OMDC DIGITAL 2026", rx+Inches(0.1), cy+Inches(0.05),
    col_w-Inches(0.2), Inches(0.35), sz=11, bold=True, color=TEAL_DARK)

visions = [
    "Booking online 24/7 via website, app & PWA",
    "Self check-in dengan Smart eKiosk tanpa antrean",
    "Mobile app lengkap dengan membership & loyalty",
    "Ekosistem digital terintegrasi lintas cabang",
    "Dashboard analytics real-time untuk manajemen",
]
for i, v in enumerate(visions):
    iy = cy + Inches(0.55) + i*Inches(0.6)
    rrect(s, rx, iy, col_w-Inches(0.1), Inches(0.52), WHITE, LIGHT_GRAY, 0.5)
    rect(s, rx, iy, Pt(4), Inches(0.52), TEAL)
    txt(s, "✓  "+v, rx+Inches(0.12), iy+Inches(0.06),
        col_w-Inches(0.3), Inches(0.42), sz=11, color=DARK_TEXT)

# Bottom KPI strip
ky = H - Inches(1.0)
rect(s, Inches(0.4), ky, W-Inches(0.8), Inches(0.5), NAVY)
txt(s, "🎯  40% Efisiensi Registrasi  •  35% Peningkatan Konversi  •  50% Pengurangan Admin Workload  •  25% Repeat Visit",
    Inches(0.6), ky+Pt(7), W-Inches(1), Inches(0.38),
    sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
print("✓ Slide 2: Executive Summary")

# ─────────────────────────────────────────
# SLIDE 3 — CURRENT CHALLENGES
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Tantangan Operasional yang Harus Diselesaikan",
            "Identifikasi pain points utama yang menghambat pertumbuhan OMDC")
footer(s)

challenges = [
    ("01", "Registrasi Manual", "Pasien mengisi formulir fisik setiap kunjungan. Proses 15–20 menit per pasien menciptakan bottleneck fatal di peak hours, menurunkan kepuasan dan kapasitas harian."),
    ("02", "Antrean Tidak Terukur", "Tanpa sistem queue digital, pasien tidak tahu estimasi waktu tunggu. Ketidakpastian ini mendorong mereka beralih ke klinik kompetitor yang lebih terorganisir."),
    ("03", "Sistem Terfragmentasi", "Data pasien, jadwal dokter, rekam medis, dan pembayaran tersebar di sistem berbeda. Tim admin kehilangan 3–4 jam per hari hanya untuk rekonsiliasi data manual."),
    ("04", "Zero Analytics", "Manajemen tidak punya visibilitas real-time atas kinerja cabang, utilisasi dokter, tren revenue, atau perilaku pasien. Keputusan bisnis dibuat tanpa data akurat."),
    ("05", "Tidak Ada Self-Service", "Pasien harus selalu berinteraksi dengan staff untuk setiap kebutuhan: pendaftaran, cek jadwal, pembayaran. Tidak ada kanal digital mandiri sama sekali."),
    ("06", "Loyalitas Tidak Terkelola", "Tidak ada program loyalitas digital. Pasien lama tidak mendapat insentif untuk kembali, dan tidak ada mekanisme re-engagement pasca-kunjungan yang efektif."),
]

cols = 3
cw = (W - Inches(0.9)) / cols
ch = (H - ht - Inches(0.55)) / 2

for i, (num, title, desc) in enumerate(challenges):
    col = i % cols; row = i // cols
    cx = Inches(0.4) + col * cw
    cy = ht + Inches(0.15) + row * ch
    rrect(s, cx+Pt(4), cy+Pt(4), cw-Inches(0.15), ch-Inches(0.15), WHITE, LIGHT_GRAY, 0.5)
    rect(s, cx+Pt(4), cy+Pt(4), Pt(4), ch-Inches(0.15), MAGENTA)
    txt(s, num, cx+Inches(0.18), cy+Inches(0.1), Inches(0.6), Inches(0.4),
        sz=11, bold=True, color=MAGENTA)
    txt(s, title, cx+Inches(0.18), cy+Inches(0.45), cw-Inches(0.5), Inches(0.4),
        sz=14, bold=True, color=DARK_TEXT)
    txt(s, desc, cx+Inches(0.18), cy+Inches(0.82), cw-Inches(0.45), ch-Inches(1.1),
        sz=10, color=GRAY)
print("✓ Slide 3: Current Challenges")

# ─────────────────────────────────────────
# SLIDE 4 — FUTURE PATIENT JOURNEY
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Future State: Perjalanan Pasien Digital OMDC",
            "Dari search hingga re-booking — setiap touchpoint didigitalisasi secara seamless")
footer(s)

steps = [
    ("🔍", "Cari", "Google Search\nSEO #1"),
    ("📅", "Booking", "App / PWA\nWhatsApp"),
    ("🖥️", "Check-In", "Smart eKiosk\nQR Scan"),
    ("🦷", "Perawatan", "Dokter +\nRekam Digital"),
    ("💳", "Pembayaran", "Digital Pay\nE-Receipt"),
    ("⭐", "Loyalty", "Poin Reward\nMembership"),
    ("🔄", "Re-Booking", "Push Notif\nReminder"),
]

sw = (W - Inches(1.0)) / len(steps)
sy = ht + Inches(0.5)
sh = H - ht - Inches(1.2)

for i, (ico, title, sub) in enumerate(steps):
    sx = Inches(0.4) + i * sw
    # Connector line
    if i < len(steps)-1:
        ary = sy + sh/2 - Pt(6)
        rect(s, sx+sw-Inches(0.05), ary, Inches(0.3), Pt(3), TEAL)
        txt(s, "›", sx+sw+Inches(0.1), ary-Inches(0.15), Inches(0.2), Inches(0.3),
            sz=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    # Card
    fc = NAVY if i==0 else (MAGENTA if i==2 else (TEAL_DARK if i==4 else WHITE))
    tc = WHITE if fc!=WHITE else DARK_TEXT
    rrect(s, sx+Pt(3), sy, sw-Pt(10), sh, fc, LIGHT_GRAY if fc==WHITE else None, 0.5 if fc==WHITE else None)
    if fc==WHITE: rect(s, sx+Pt(3), sy, Pt(4), sh, MAGENTA)
    txt(s, ico, sx+Pt(3), sy+Inches(0.2), sw-Pt(10), Inches(0.5),
        sz=22, align=PP_ALIGN.CENTER, color=WHITE if fc!=WHITE else MAGENTA)
    txt(s, title, sx+Pt(3), sy+Inches(0.75), sw-Pt(10), Inches(0.4),
        sz=13, bold=True, color=tc, align=PP_ALIGN.CENTER)
    txt(s, sub, sx+Pt(3), sy+Inches(1.18), sw-Pt(10), sh-Inches(1.4),
        sz=9, color=tc if fc!=WHITE else GRAY, align=PP_ALIGN.CENTER)
print("✓ Slide 4: Future Patient Journey")

# ─────────────────────────────────────────
# SLIDE 5 — DIGITAL ECOSYSTEM
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Ekosistem Digital OMDC: Satu Platform Terintegrasi",
            "Arsitektur berlapis yang menghubungkan setiap stakeholder dalam satu ekosistem cerdas")
footer(s)

layers = [
    (NAVY,    WHITE,   "PASIEN",                "Booking • Check-In • Pembayaran • Loyalty"),
    (MAGENTA, WHITE,   "WEBSITE  +  PWA  +  MOBILE APP", "Akuisisi • Engagement • Self-Service"),
    (RGBColor(0x1A,0x6B,0x8A), WHITE, "SMART eKIOSK",  "Registrasi • Antrian • Pembayaran • QR Scan"),
    (TEAL_DARK, WHITE, "OPERASIONAL KLINIK",    "Rekam Medis • Jadwal Dokter • Billing"),
    (RGBColor(0x2D,0x3F,0x7A), WHITE, "MANAGEMENT DASHBOARD", "Analytics • BI • Multi-Cabang • Reporting"),
]

total_h = H - ht - Inches(0.55)
lh = total_h / len(layers)
max_w = W - Inches(1.6)

for i, (fc, tc, lbl, desc) in enumerate(layers):
    lw = max_w - i * Inches(0.7)
    lx = Inches(0.8) + i * Inches(0.35)
    ly = ht + Inches(0.1) + i * lh
    rrect(s, lx, ly, lw, lh-Pt(6), fc)
    txt(s, lbl, lx+Inches(0.3), ly+Inches(0.08), lw/2, lh-Pt(12),
        sz=13, bold=True, color=tc, align=PP_ALIGN.LEFT)
    txt(s, desc, lx+lw/2, ly+Inches(0.08), lw/2-Inches(0.3), lh-Pt(12),
        sz=10, color=RGBColor(0xCC,0xDD,0xEE) if fc!=WHITE else GRAY,
        align=PP_ALIGN.RIGHT)
    # Step number
    txt(s, str(i+1), W-Inches(0.6), ly+Inches(0.08), Inches(0.4), lh-Pt(12),
        sz=18, bold=True, color=RGBColor(0x3D,0x5A,0x9A), align=PP_ALIGN.CENTER)
print("✓ Slide 5: Digital Ecosystem")

# ─────────────────────────────────────────
# SLIDE 6 — WEBSITE REVAMP
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Website Revamp: Pintu Digital Utama OMDC",
            "Dari website informasional menjadi mesin akuisisi pasien yang powerful")
footer(s)

# Left: Features
lw = W * 0.52
cy = ht + Inches(0.2)

features = [
    ("🔖", "Direktori Dokter & Layanan", "Profil dokter dengan spesialisasi, jadwal & ulasan. Katalog lengkap layanan gigi dengan estimasi biaya."),
    ("📍", "Booking Online & Branch Locator", "Reservasi real-time dengan kalender dokter. Peta interaktif semua cabang OMDC."),
    ("💬", "Pusat Promo & WhatsApp API", "Landing page promo dinamis untuk campaign seasonal. Integrasi WhatsApp Business untuk follow-up otomatis."),
    ("⚡", "SEO-Optimized & Fast Loading", "Technical SEO terbaik, loading < 2 detik, mobile-first responsive, Core Web Vitals sempurna."),
]

fh = (H - ht - Inches(0.55)) / len(features)
for i, (ico, ttl, dsc) in enumerate(features):
    fy = cy + i * fh
    rrect(s, Inches(0.4), fy+Pt(4), lw-Inches(0.2), fh-Pt(8), WHITE, LIGHT_GRAY, 0.5)
    rect(s, Inches(0.4), fy+Pt(4), Pt(4), fh-Pt(8), MAGENTA)
    txt(s, ico, Inches(0.55), fy+Inches(0.12), Inches(0.5), Inches(0.45),
        sz=18, align=PP_ALIGN.CENTER, color=MAGENTA)
    txt(s, ttl, Inches(1.1), fy+Inches(0.1), lw-Inches(1.3), Inches(0.38),
        sz=13, bold=True, color=DARK_TEXT)
    txt(s, dsc, Inches(1.1), fy+Inches(0.48), lw-Inches(1.3), fh-Inches(0.65),
        sz=10, color=GRAY)

# Right: Stats
rx = lw + Inches(0.3)
rw = W - lw - Inches(0.5)

rrect(s, rx, cy, rw, Inches(1.0), NAVY)
txt(s, "EXPECTED RESULTS", rx+Inches(0.15), cy+Inches(0.1), rw-Inches(0.3), Inches(0.35),
    sz=10, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)
txt(s, "+35%", rx, cy+Inches(0.35), rw, Inches(0.5),
    sz=30, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(s, "Peningkatan Konversi", rx, cy+Inches(0.78), rw, Inches(0.25),
    sz=10, color=WHITE, align=PP_ALIGN.CENTER)

rrect(s, rx, cy+Inches(1.1), rw, Inches(1.0), MAGENTA)
txt(s, "+40%", rx, cy+Inches(1.2), rw, Inches(0.55),
    sz=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Permintaan Appointment", rx, cy+Inches(1.72), rw, Inches(0.25),
    sz=10, color=WHITE, align=PP_ALIGN.CENTER)

# Mock browser frame
bfy = cy + Inches(2.25)
bfh = H - bfy - Inches(0.5)
rrect(s, rx, bfy, rw, bfh, WHITE, LIGHT_GRAY, 1)
rect(s, rx, bfy, rw, Inches(0.32), RGBColor(0xF0,0xF0,0xF3))
rect(s, rx+Inches(0.08), bfy+Pt(8), Inches(0.12), Inches(0.12), RGBColor(0xFF,0x6B,0x6B))
rect(s, rx+Inches(0.26), bfy+Pt(8), Inches(0.12), Inches(0.12), RGBColor(0xFF,0xD9,0x3D))
rect(s, rx+Inches(0.44), bfy+Pt(8), Inches(0.12), Inches(0.12), RGBColor(0x6B,0xCB,0x6B))
rrect(s, rx+Inches(0.65), bfy+Pt(6), rw-Inches(0.8), Inches(0.18), LIGHT_GRAY)
txt(s, "omdc.co.id", rx+Inches(0.65), bfy+Pt(6), rw-Inches(0.8), Inches(0.18),
    sz=8, color=GRAY, align=PP_ALIGN.CENTER)
rect(s, rx, bfy+Inches(0.32), rw, Inches(0.45), MAGENTA)
txt(s, "OMDC — Klinik Gigi Terpercaya Indonesia", rx, bfy+Inches(0.4), rw, Inches(0.32),
    sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for j, lnk in enumerate(["Dokter", "Layanan", "Cabang", "Promo", "Booking"]):
    lx2 = rx + Inches(0.1) + j*(rw/5)
    txt(s, lnk, lx2, bfy+Inches(0.82), rw/5-Pt(4), Inches(0.25),
        sz=8, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)
print("✓ Slide 6: Website Revamp")

# ─────────────────────────────────────────
# SLIDE 7 — SEO & SALES BOOSTER
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "SEO & Strategi Dominasi Digital",
            "Menguasai halaman pertama Google untuk semua kata kunci dental di Indonesia")
footer(s)

# Left: keyword targets
kws = ["Klinik Gigi Jakarta", "Dokter Gigi Terbaik", "Scaling Gigi",
       "Behel Gigi Jakarta", "Veneer Gigi", "Implant Gigi", "Perawatan Saluran Akar"]
lw2 = W * 0.38
cy2 = ht + Inches(0.2)
txt(s, "TARGET KEYWORDS", Inches(0.4), cy2, lw2, Inches(0.38),
    sz=11, bold=True, color=MAGENTA)

for i, kw in enumerate(kws):
    ky2 = cy2 + Inches(0.45) + i * Inches(0.57)
    rrect(s, Inches(0.4), ky2, lw2-Inches(0.2), Inches(0.45),
          RGBColor(0xFF,0xF0,0xF8), LIGHT_GRAY, 0.5)
    txt(s, "🔍", Inches(0.5), ky2+Pt(4), Inches(0.3), Inches(0.35), sz=12)
    txt(s, kw, Inches(0.85), ky2+Pt(5), lw2-Inches(0.7), Inches(0.35),
        sz=12, color=DARK_TEXT)
    # Volume bar
    bar_w = (lw2-Inches(1.2)) * (0.95 - i*0.1)
    rect(s, Inches(0.4)+(lw2-Inches(0.2))-bar_w-Inches(0.05), ky2+Inches(0.3),
         bar_w, Pt(4), MAGENTA)

# Center: growth arrow
mx = lw2 + Inches(0.2)
mw = Inches(2.5)
txt(s, "📈", mx, ht+Inches(1.5), mw, Inches(1.2), sz=60, align=PP_ALIGN.CENTER)
txt(s, "+100%", mx, ht+Inches(2.7), mw, Inches(0.75),
    sz=32, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)
txt(s, "Organic Growth Target", mx, ht+Inches(3.45), mw, Inches(0.4),
    sz=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# Right: approach cards
rx3 = lw2 + mw + Inches(0.4)
rw3 = W - rx3 - Inches(0.4)
approaches = [
    ("Technical SEO", "Site speed, Core Web Vitals, schema markup, crawlability"),
    ("Content Marketing", "Artikel blog edukasi gigi, FAQ, panduan perawatan"),
    ("Local SEO", "Google Business tiap cabang, local citations, review management"),
    ("Link Building", "Backlink authority dari situs kesehatan & media nasional"),
]
for i, (ttl, dsc) in enumerate(approaches):
    ay = ht + Inches(0.2) + i * Inches(1.2)
    rrect(s, rx3, ay, rw3, Inches(1.1), WHITE, LIGHT_GRAY, 0.5)
    rect(s, rx3, ay, Pt(4), Inches(1.1), TEAL)
    txt(s, ttl, rx3+Inches(0.15), ay+Inches(0.1), rw3-Inches(0.2), Inches(0.38),
        sz=12, bold=True, color=DARK_TEXT)
    txt(s, dsc, rx3+Inches(0.15), ay+Inches(0.48), rw3-Inches(0.2), Inches(0.55),
        sz=10, color=GRAY)
print("✓ Slide 7: SEO & Sales Booster")

# ─────────────────────────────────────────
# SLIDE 8 — PWA PLATFORM
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Progressive Web App (PWA): Pengalaman App Tanpa Instalasi",
            "Akses semua fitur via browser smartphone — tanpa download, tanpa hambatan")
footer(s)

feats = [
    ("📅", "Booking Appointment", "Reservasi real-time dengan pilihan dokter, jadwal, dan cabang. Konfirmasi instan via push notifikasi dan WhatsApp. Reminder otomatis H-1 untuk reduce no-show."),
    ("🏆", "Membership & Loyalty", "Program poin digital: earn setiap kunjungan, redeem untuk diskon perawatan. Tier keanggotaan Silver, Gold, Platinum dengan benefit eksklusif."),
    ("📋", "Riwayat Perawatan", "Akses digital seluruh rekam medis, foto before-after, rencana perawatan, dan jadwal kontrol yang direkomendasikan dokter kapan saja."),
    ("💳", "Pembayaran Digital", "Ekosistem payment lengkap: transfer bank, GoPay, OVO, DANA, kartu kredit/debit, QRIS. Riwayat transaksi dan e-receipt dapat diunduh kapan saja."),
]

cw8 = (W - Inches(1.0)) / 2
ch8 = (H - ht - Inches(0.6)) / 2

for i, (ico, ttl, dsc) in enumerate(feats):
    col = i % 2; row = i // 2
    fx = Inches(0.4) + col * (cw8 + Inches(0.2))
    fy = ht + Inches(0.15) + row * (ch8 + Inches(0.1))
    rrect(s, fx, fy, cw8, ch8, WHITE, LIGHT_GRAY, 0.5)
    rect(s, fx, fy, cw8, Pt(4), MAGENTA)
    # Icon circle
    oval(s, fx+Inches(0.2), fy+Inches(0.1), Inches(0.55), Inches(0.55), LIGHT_MAG)
    txt(s, ico, fx+Inches(0.2), fy+Inches(0.1), Inches(0.55), Inches(0.55),
        sz=20, align=PP_ALIGN.CENTER)
    txt(s, ttl, fx+Inches(0.9), fy+Inches(0.12), cw8-Inches(1.0), Inches(0.4),
        sz=13, bold=True, color=DARK_TEXT)
    txt(s, dsc, fx+Inches(0.2), fy+Inches(0.72), cw8-Inches(0.4), ch8-Inches(0.9),
        sz=10, color=GRAY)
print("✓ Slide 8: PWA Platform")

# ─────────────────────────────────────────
# SLIDE 9 — ANDROID & iOS APPS
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Aplikasi Native Android & iOS",
            "Pengalaman mobile premium yang dirancang untuk standar klinik gigi kelas dunia")
footer(s)

# 3 phone mockups
phone_w = Inches(2.5)
phone_h = Inches(4.8)
phone_y = ht + Inches(0.25)
screens = [
    (NAVY, "Dashboard", ["Selamat datang, Budi 👋", "", "Appointment Berikutnya", "Rabu 25 Jun — 10.00", "Dr. Sari Dewi • Cabang Sudirman", "", "💎 Gold Member", "1.250 poin aktif"]),
    (MAGENTA, "Booking", ["Pilih Layanan", "[ Scaling & Polishing ]", "[ Behel Gigi ]", "[ Veneer ]", "", "Pilih Dokter", "▶ Dr. Rina ● Dr. Andi ● Dr. Sari", "", "📅 Pilih Tanggal & Waktu", "[  25 Jun  ] [ 10.00 ]", "", "[ Konfirmasi Booking ]"]),
    (TEAL_DARK, "Rewards", ["💎 Gold Member", "━━━━━━━━━━━━━━", "1.250 / 2.000 Poin", "[████████░░] 63%", "", "Reward Tersedia:", "• Diskon 15% Scaling", "• Free Konsultasi", "• Priority Queue"]),
]
for i, (fc, ttl, lines) in enumerate(screens):
    px = Inches(0.5) + i * Inches(4.28)
    # Phone frame
    rrect(s, px, phone_y, phone_w, phone_h, RGBColor(0xE8,0xE8,0xEE), None, None)
    rrect(s, px+Pt(8), phone_y+Pt(8), phone_w-Pt(16), phone_h-Pt(16), fc)
    # Status bar
    txt(s, "9:41  ●●●●  WiFi  🔋", px+Pt(12), phone_y+Pt(12),
        phone_w-Pt(24), Pt(16), sz=6, color=WHITE, align=PP_ALIGN.CENTER)
    # App header
    rect(s, px+Pt(8), phone_y+Pt(28), phone_w-Pt(16), Pt(28), RGBColor(0x0,0x0,0x0,))
    txt(s, f"OMDC — {ttl}", px+Pt(12), phone_y+Pt(29),
        phone_w-Pt(24), Pt(22), sz=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Content
    for j, line in enumerate(lines):
        txt(s, line, px+Pt(14), phone_y+Pt(62)+j*Pt(22),
            phone_w-Pt(28), Pt(20), sz=8, color=WHITE if line else WHITE,
            align=PP_ALIGN.LEFT)
    # Screen label
    rrect(s, px, phone_y+phone_h+Inches(0.1), phone_w, Inches(0.35), fc)
    txt(s, f"Screen {i+1}: {ttl}", px, phone_y+phone_h+Inches(0.1),
        phone_w, Inches(0.35), sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Right panel: feature highlights
rp_x = Inches(0.5) + 3*Inches(4.28) - Inches(0.3)
# ... adjust for 3 phones
rp_x = Inches(13.0) - Inches(0.4)

# Description bullets on right side under phones
desc_items = [
    "React Native — satu codebase untuk Android & iOS",
    "Push notifications cerdas & contextual",
    "Offline mode untuk data dasar tersedia",
    "Face ID / Fingerprint authentication",
    "App Store & Google Play distribution",
]
for i, d in enumerate(desc_items):
    dy = ht + Inches(0.3) + i * Inches(0.55)
    rrect(s, W-Inches(2.8), dy, Inches(2.4), Inches(0.45), WHITE, LIGHT_GRAY, 0.5)
    rect(s, W-Inches(2.8), dy, Pt(4), Inches(0.45), TEAL)
    txt(s, d, W-Inches(2.68), dy+Pt(5), Inches(2.25), Inches(0.38), sz=9, color=DARK_TEXT)
print("✓ Slide 9: Android & iOS Apps")

# ─────────────────────────────────────────
# SLIDE 10 — SMART eKIOSK PLATFORM
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Smart eKiosk: Revolusi Layanan Mandiri Pasien",
            "Terminal self-service yang mengurangi beban front desk hingga 50% dan mempercepat alur layanan")
footer(s)

phases = [
    ("🚶", "Kedatangan", "Pasien tiba di klinik\nMenuju terminal eKiosk\nIdentifikasi diri tersedia"),
    ("📲", "QR Scan", "Scan QR dari aplikasi\natau input NIK manual\nVerifikasi appointment"),
    ("✅", "Check-In", "Konfirmasi data diri\nCetak/kirim nomor antrian\nNotif ke smartphone"),
    ("🦷", "Perawatan", "Panggilan antrian digital\nMasuk ruang perawatan\nDokter akses EMR"),
    ("💳", "Pembayaran", "Tagihan otomatis muncul\nPilih metode payment\nCetak / kirim e-receipt"),
    ("⭐", "Selesai", "Rating & feedback\nUpdate poin loyalty\nReminder kontrol berikutnya"),
]

pw = (W - Inches(1.0)) / len(phases)
py_base = ht + Inches(0.3)
ph = H - ht - Inches(1.0)

for i, (ico, ttl, dsc) in enumerate(phases):
    px2 = Inches(0.4) + i * pw
    fc = [NAVY, MAGENTA, TEAL_DARK, NAVY2, MAGENTA, TEAL_DARK][i]
    # Flow card
    rrect(s, px2+Pt(3), py_base, pw-Pt(6), ph, fc)
    # Step number
    oval(s, px2+Inches(0.1), py_base+Inches(0.1), Inches(0.45), Inches(0.45),
         RGBColor(0xFF,0xFF,0xFF) if fc!=WHITE else MAGENTA)
    txt(s, str(i+1), px2+Inches(0.1), py_base+Inches(0.1), Inches(0.45), Inches(0.45),
        sz=12, bold=True, color=fc, align=PP_ALIGN.CENTER)
    txt(s, ico, px2, py_base+Inches(0.68), pw, Inches(0.6),
        sz=28, align=PP_ALIGN.CENTER, color=WHITE)
    txt(s, ttl, px2, py_base+Inches(1.35), pw, Inches(0.38),
        sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, dsc, px2+Pt(6), py_base+Inches(1.78), pw-Pt(12), ph-Inches(2.0),
        sz=9, color=RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER)
    # Arrow connector
    if i < len(phases)-1:
        ax = px2+pw-Pt(4)
        ay = py_base+ph/2-Pt(10)
        rect(s, ax, ay, Pt(6), Pt(3), WHITE)
        txt(s, "›", ax-Pt(2), ay-Pt(10), Pt(18), Pt(22),
            sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
print("✓ Slide 10: Smart eKiosk Platform")

# ─────────────────────────────────────────
# SLIDE 11 — eKIOSK FEATURES
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Fitur Unggulan Platform Smart eKiosk",
            "Teknologi enterprise-grade dalam form factor yang elegan dan mudah digunakan")
footer(s)

feats11 = [
    ("🪪", "Registrasi Mandiri", "Pendaftaran pasien baru dalam 2 menit. Verifikasi NIK, QR membership, atau kartu pasien. Integrasi langsung ke database OMDC terpusat."),
    ("🔢", "Queue Management", "Antrian digital cerdas dengan estimasi waktu akurat. Real-time queue display. Prioritas otomatis untuk lansia & emergency."),
    ("💰", "Multi-Payment", "Tunai (kembalian otomatis), semua kartu debit/kredit, e-wallet populer, QRIS universal. Integrasi akuntansi langsung."),
    ("🧾", "Print & Digital Receipt", "Thermal print instan atau kirim ke email/WhatsApp. Riwayat transaksi tersimpan di akun pasien permanent."),
    ("📡", "Real-Time Sync", "Sinkronisasi zero-delay ke semua sistem: EMR, billing, jadwal. Event-driven architecture yang handal dan fault-tolerant."),
    ("📊", "Monitoring Dashboard", "Dashboard operasional kiosk real-time: status online, volume transaksi, error alert, dan laporan performa harian."),
]

cw11 = (W - Inches(1.0)) / 3
ch11 = (H - ht - Inches(0.6)) / 2

for i, (ico, ttl, dsc) in enumerate(feats11):
    col = i % 3; row = i // 3
    fx = Inches(0.4) + col * (cw11 + Inches(0.1))
    fy = ht + Inches(0.15) + row * (ch11 + Inches(0.1))
    rrect(s, fx, fy, cw11, ch11, WHITE, LIGHT_GRAY, 0.5)
    ac = [MAGENTA, TEAL, NAVY, TEAL, MAGENTA, NAVY][i]
    rect(s, fx, fy, cw11, Pt(4), ac)
    oval(s, fx+Inches(0.15), fy+Inches(0.12), Inches(0.55), Inches(0.55),
         RGBColor(0xFC,0xE7,0xF3) if ac==MAGENTA else RGBColor(0xE0,0xFB,0xFC))
    txt(s, ico, fx+Inches(0.15), fy+Inches(0.12), Inches(0.55), Inches(0.55),
        sz=20, align=PP_ALIGN.CENTER)
    txt(s, ttl, fx+Inches(0.85), fy+Inches(0.14), cw11-Inches(1.0), Inches(0.38),
        sz=13, bold=True, color=DARK_TEXT)
    txt(s, dsc, fx+Inches(0.15), fy+Inches(0.73), cw11-Inches(0.3), ch11-Inches(0.9),
        sz=10, color=GRAY)
print("✓ Slide 11: eKiosk Features")

# ─────────────────────────────────────────
# SLIDE 12 — HARDWARE SPECIFICATION
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Spesifikasi Hardware eKiosk: Enterprise Grade",
            "Komponen premium untuk keandalan, estetika, dan performa optimal di lingkungan klinik 24/7")
footer(s)

# Left: kiosk diagram
kiosk_x = Inches(0.4)
kiosk_y = ht + Inches(0.25)
kiosk_w = Inches(3.5)
kiosk_h = H - ht - Inches(0.7)

# Stand
rect(s, kiosk_x+kiosk_w/2-Inches(0.35), kiosk_y+kiosk_h-Inches(0.6),
     Inches(0.7), Inches(0.6), RGBColor(0x9C,0xA3,0xAF))
# Base
rrect(s, kiosk_x+Inches(0.3), kiosk_y+kiosk_h-Inches(0.1),
      kiosk_w-Inches(0.6), Inches(0.08), RGBColor(0x6B,0x72,0x80))
# Body
rrect(s, kiosk_x+Inches(0.15), kiosk_y, kiosk_w-Inches(0.3), kiosk_h-Inches(0.65), NAVY)
# Screen
rrect(s, kiosk_x+Inches(0.35), kiosk_y+Inches(0.15),
      kiosk_w-Inches(0.7), kiosk_h-Inches(1.2), NAVY2)
# Screen content
txt(s, "OMDC", kiosk_x+Inches(0.4), kiosk_y+Inches(0.3), kiosk_w-Inches(0.8), Inches(0.5),
    sz=18, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)
txt(s, "Selamat Datang", kiosk_x+Inches(0.4), kiosk_y+Inches(0.85), kiosk_w-Inches(0.8), Inches(0.4),
    sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rrect(s, kiosk_x+Inches(0.65), kiosk_y+Inches(1.4), kiosk_w-Inches(1.3), Inches(0.45), MAGENTA)
txt(s, "[ SCAN QR ]", kiosk_x+Inches(0.65), kiosk_y+Inches(1.4), kiosk_w-Inches(1.3), Inches(0.45),
    sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rrect(s, kiosk_x+Inches(0.65), kiosk_y+Inches(2.0), kiosk_w-Inches(1.3), Inches(0.45), TEAL_DARK)
txt(s, "[ DAFTAR BARU ]", kiosk_x+Inches(0.65), kiosk_y+Inches(2.0), kiosk_w-Inches(1.3), Inches(0.45),
    sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# Printer slot
rect(s, kiosk_x+Inches(0.8), kiosk_y+kiosk_h-Inches(1.15),
     kiosk_w-Inches(1.6), Inches(0.1), RGBColor(0x2D,0x3F,0x7A))
# Card reader
rrect(s, kiosk_x+Inches(0.65), kiosk_y+kiosk_h-Inches(0.95),
      Inches(0.6), Inches(0.25), RGBColor(0x2D,0x3F,0x7A))
txt(s, "Smart eKiosk™", kiosk_x+Inches(0.15), kiosk_y+kiosk_h-Inches(0.62),
    kiosk_w-Inches(0.3), Inches(0.25), sz=8, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)

# Right: specs table
rx12 = kiosk_x + kiosk_w + Inches(0.4)
rw12 = W - rx12 - Inches(0.4)

specs = [
    ("LAYAR", [("Ukuran", "32\" Full HD IPS Touchscreen"), ("Touch", "10-point multi-touch kapasitif"), ("Brightness", "400 nits, anti-glare"), ("Body", "Stainless steel premium")]),
    ("COMPUTING", [("Processor", "Intel Core i5 Gen 12 Industrial"), ("RAM", "16GB DDR4"), ("Storage", "256GB NVMe SSD, enkripsi HW"), ("OS", "Windows 10 IoT Enterprise LTSC")]),
    ("PERIPHERAL", [("Printer", "Thermal 80mm, 200mm/s"), ("Scanner", "2D QR & barcode omnidirectional"), ("Card Reader", "EMV chip + contactless NFC"), ("Power Backup", "Built-in UPS 2 jam")]),
    ("KONEKTIVITAS", [("Network", "Gigabit Ethernet + WiFi 6"), ("Security", "Kunci fisik & alarm anti-tamper"), ("Remote Mgmt", "24/7 monitoring & management"), ("Garansi", "2 tahun onsite replacement")]),
]

sh12 = (H - ht - Inches(0.6)) / len(specs)
for i, (cat, items) in enumerate(specs):
    sy12 = ht + Inches(0.15) + i * sh12
    rrect(s, rx12, sy12+Pt(3), rw12, sh12-Pt(6), WHITE, LIGHT_GRAY, 0.5)
    rect(s, rx12, sy12+Pt(3), Pt(4), sh12-Pt(6), [MAGENTA,TEAL,NAVY,TEAL_DARK][i])
    txt(s, cat, rx12+Inches(0.15), sy12+Inches(0.06), rw12*0.35, Inches(0.3),
        sz=10, bold=True, color=[MAGENTA,TEAL,NAVY,TEAL_DARK][i])
    for j, (k, v) in enumerate(items):
        col_x = rx12 + Inches(0.15) + (j//2) * (rw12/2)
        row_y = sy12 + Inches(0.38) + (j%2) * Inches(0.3)
        txt(s, f"{k}: ", col_x, row_y, Inches(1.1), Inches(0.28),
            sz=9, bold=True, color=DARK_TEXT)
        txt(s, v, col_x+Inches(1.05), row_y, rw12/2-Inches(1.1), Inches(0.28),
            sz=9, color=GRAY)
print("✓ Slide 12: Hardware Specification")

# ─────────────────────────────────────────
# SLIDE 13 — TECHNOLOGY ARCHITECTURE
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Arsitektur Teknologi: Enterprise-Grade Stack",
            "Fondasi teknologi modern yang scalable, secure, dan siap tumbuh bersama OMDC selama 5–10 tahun")
footer(s)

layers13 = [
    (RGBColor(0x0D,0x1B,0x3E), "☁️  INFRASTRUCTURE", "AWS Cloud: EC2 • RDS Aurora • S3 • CloudFront CDN • Auto-scaling • 99.9% SLA • Multi-region DR"),
    (RGBColor(0x15,0x28,0x55), "🗄️  DATABASE & CACHE", "PostgreSQL (primary + read replicas) • Redis caching • Enkripsi AES-256 • Backup tiap 6 jam • Replikasi real-time"),
    (RGBColor(0x1A,0x47,0x8A), "⚙️  BACKEND SERVICES", "NestJS Microservices: Auth • Booking • EMR • Payment • Notification • Analytics • RabbitMQ event bus"),
    (RGBColor(0x00,0x6D,0x77), "💻  FRONTEND & MOBILE", "Next.js 14 App Router (Website + PWA) • React Native (Android + iOS) • TypeScript end-to-end"),
    (RGBColor(0x00,0x96,0x9B), "🔌  INTEGRATIONS", "Midtrans & Xendit (Payment) • Firebase FCM (Notif) • WhatsApp Business API • Google Analytics 4 • OAuth 2.0"),
]

lh13 = (H - ht - Inches(0.55)) / len(layers13)
for i, (fc, lbl, tech) in enumerate(layers13):
    lw13 = W - Inches(0.8) - i * Inches(0.5)
    lx13 = Inches(0.4) + i * Inches(0.25)
    ly13 = ht + Inches(0.1) + i * lh13
    rrect(s, lx13, ly13+Pt(2), lw13, lh13-Pt(4), fc)
    txt(s, lbl, lx13+Inches(0.25), ly13+Inches(0.1), Inches(3.5), lh13-Pt(8),
        sz=12, bold=True, color=WHITE)
    txt(s, tech, lx13+Inches(3.8), ly13+Inches(0.1), lw13-Inches(4.0), lh13-Pt(8),
        sz=10, color=RGBColor(0xB0,0xCC,0xE8))
    txt(s, str(i+1), W-Inches(0.55), ly13+Inches(0.1), Inches(0.35), lh13-Pt(8),
        sz=14, bold=True, color=RGBColor(0x4A,0x6A,0xAA), align=PP_ALIGN.CENTER)
print("✓ Slide 13: Technology Architecture")

# ─────────────────────────────────────────
# SLIDE 14 — SECURITY & COMPLIANCE
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Keamanan Data & Kepatuhan Enterprise",
            "Framework keamanan berlapis yang melindungi data sensitif pasien dan kepercayaan bisnis OMDC")
footer(s)

secs = [
    ("🔒", "Enkripsi End-to-End", "AES-256 at-rest, TLS 1.3 in-transit", "Semua data pasien dienkripsi penuh. Kunci enkripsi dikelola AWS KMS. Zero-knowledge architecture—tidak ada akses tidak terotorisasi bahkan oleh administrator internal."),
    ("🛡️", "Access Control & IAM", "RBAC granular + MFA mandatory", "Role-based access control: admin, dokter, perawat, kasir, manajer. MFA wajib semua akun. Auto-logout sesi idle. Anomaly detection untuk akses mencurigakan."),
    ("📋", "Audit Log & Compliance", "Pencatatan forensik immutable", "Setiap akses dan modifikasi data dicatat secara permanen dan tidak bisa diubah. Kepatuhan regulasi kesehatan Indonesia, PDPA, dan standar ISO 27001 healthcare."),
    ("☁️", "Backup & Disaster Recovery", "RTO < 4 jam, RPO < 1 jam", "Backup otomatis tiap 6 jam ke 3 region AWS berbeda. Disaster recovery drill rutin per kuartal. 99.9% uptime SLA dengan arsitektur high-availability."),
]

cw14 = (W - Inches(1.0)) / 2
ch14 = (H - ht - Inches(0.6)) / 2

for i, (ico, ttl, sub, dsc) in enumerate(secs):
    col = i % 2; row = i // 2
    fx14 = Inches(0.4) + col * (cw14 + Inches(0.2))
    fy14 = ht + Inches(0.15) + row * (ch14 + Inches(0.1))
    rrect(s, fx14, fy14, cw14, ch14, WHITE, LIGHT_GRAY, 0.5)
    ac14 = [NAVY, MAGENTA, TEAL_DARK, TEAL][i]
    rect(s, fx14, fy14, cw14, Pt(4), ac14)
    # Icon badge
    rrect(s, fx14+Inches(0.2), fy14+Inches(0.12), Inches(0.6), Inches(0.6),
          RGBColor(0xEA,0xF4,0xFF) if ac14==NAVY else LIGHT_MAG)
    txt(s, ico, fx14+Inches(0.2), fy14+Inches(0.12), Inches(0.6), Inches(0.6),
        sz=22, align=PP_ALIGN.CENTER)
    txt(s, ttl, fx14+Inches(0.95), fy14+Inches(0.14), cw14-Inches(1.1), Inches(0.38),
        sz=13, bold=True, color=DARK_TEXT)
    tag(s, sub, fx14+Inches(0.95), fy14+Inches(0.52), Inches(2.5), Inches(0.26),
        RGBColor(0xEA,0xF4,0xFF) if ac14==NAVY else LIGHT_MAG,
        ac14)
    txt(s, dsc, fx14+Inches(0.2), fy14+Inches(0.88), cw14-Inches(0.4), ch14-Inches(1.05),
        sz=10, color=GRAY)
print("✓ Slide 14: Security & Compliance")

# ─────────────────────────────────────────
# SLIDE 15 — MANAGEMENT DASHBOARD
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Executive Management Dashboard",
            "Visibilitas penuh atas performa seluruh ekosistem OMDC dalam satu layar real-time")
footer(s)

# Mock dashboard frame
dash_x, dash_y = Inches(0.4), ht+Inches(0.15)
dash_w, dash_h = W-Inches(0.8), H-ht-Inches(0.6)
rrect(s, dash_x, dash_y, dash_w, dash_h, WHITE, LIGHT_GRAY, 1)

# Dashboard header bar
rect(s, dash_x, dash_y, dash_w, Inches(0.45), NAVY)
txt(s, "🏥  OMDC Command Center — Live Dashboard  |  Thu 19 Jun 2026  09:41 WIB  ●  LIVE",
    dash_x+Inches(0.2), dash_y+Pt(7), dash_w-Inches(0.4), Inches(0.32),
    sz=10, bold=True, color=WHITE)

# KPI row
kpis = [
    ("487", "Total Pasien Hari Ini", MAGENTA),
    ("Rp 94.2 Jt", "Revenue Hari Ini", TEAL_DARK),
    ("87%", "Utilisasi Dokter", NAVY),
    ("4.8 / 5.0", "Rating Pasien", RGBColor(0xF5,0x9E,0x0B)),
    ("23 Mnt", "Avg. Waktu Tunggu", RGBColor(0x10,0xB9,0x81)),
]
kpi_w = dash_w / len(kpis)
for i, (val, lbl, fc) in enumerate(kpis):
    kx = dash_x + i * kpi_w
    ky = dash_y + Inches(0.48)
    rrect(s, kx+Pt(3), ky, kpi_w-Pt(6), Inches(0.9), WHITE, LIGHT_GRAY, 0.5)
    rect(s, kx+Pt(3), ky, kpi_w-Pt(6), Pt(3), fc)
    txt(s, val, kx+Pt(3), ky+Inches(0.08), kpi_w-Pt(6), Inches(0.5),
        sz=18, bold=True, color=fc, align=PP_ALIGN.CENTER)
    txt(s, lbl, kx+Pt(3), ky+Inches(0.55), kpi_w-Pt(6), Inches(0.32),
        sz=8, color=GRAY, align=PP_ALIGN.CENTER)

# Charts area
chart_y = dash_y + Inches(1.45)
chart_h = dash_h - Inches(1.5)

# Left chart: Revenue trend bars
chart_lx = dash_x + Inches(0.15)
chart_lw = dash_w * 0.38
rrect(s, chart_lx, chart_y, chart_lw, chart_h, WHITE, LIGHT_GRAY, 0.5)
txt(s, "Revenue Trend — 6 Bulan", chart_lx+Inches(0.1), chart_y+Inches(0.08),
    chart_lw-Inches(0.2), Inches(0.3), sz=10, bold=True, color=DARK_TEXT)
months = ["Jan","Feb","Mar","Apr","Mei","Jun"]
vals = [62, 71, 68, 78, 85, 94]
bar_area_h = chart_h - Inches(0.65)
bar_w2 = (chart_lw - Inches(0.4)) / len(months)
for j, (m, v) in enumerate(zip(months, vals)):
    bh2 = bar_area_h * v / 100
    bx2 = chart_lx + Inches(0.2) + j * bar_w2
    by2 = chart_y + Inches(0.42) + bar_area_h - bh2
    rrect(s, bx2+Pt(3), by2, bar_w2-Pt(6), bh2, MAGENTA)
    txt(s, m, bx2, chart_y+Inches(0.42)+bar_area_h+Pt(3), bar_w2, Pt(14), sz=7, color=GRAY, align=PP_ALIGN.CENTER)
    txt(s, f"{v}%", bx2, by2-Pt(14), bar_w2, Pt(12), sz=7, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)

# Middle chart: Appointment funnel
chart_mx = chart_lx + chart_lw + Inches(0.15)
chart_mw = dash_w * 0.3
rrect(s, chart_mx, chart_y, chart_mw, chart_h, WHITE, LIGHT_GRAY, 0.5)
txt(s, "Booking Funnel", chart_mx+Inches(0.1), chart_y+Inches(0.08),
    chart_mw-Inches(0.2), Inches(0.3), sz=10, bold=True, color=DARK_TEXT)
funnel = [("Website Visit", 100, LIGHT_GRAY), ("Lihat Layanan", 72, RGBColor(0xFC,0xE7,0xF3)),
          ("Mulai Booking", 45, LIGHT_MAG), ("Confirmed", 35, MAGENTA)]
for j, (lbl, pct, fc) in enumerate(funnel):
    fw = (chart_mw-Inches(0.3)) * pct/100
    fx2 = chart_mx + Inches(0.15) + (chart_mw-Inches(0.3)-fw)/2
    fy2 = chart_y + Inches(0.5) + j * Inches(0.55)
    rrect(s, fx2, fy2, fw, Inches(0.42), fc)
    txt(s, f"{lbl}  {pct}%", chart_mx+Inches(0.1), fy2+Inches(0.08), chart_mw-Inches(0.2), Inches(0.28),
        sz=8, color=DARK_TEXT if fc!=MAGENTA else WHITE, align=PP_ALIGN.CENTER)

# Right: Branch performance
chart_rx = chart_mx + chart_mw + Inches(0.15)
chart_rw = dash_w - (chart_rx - dash_x) - Inches(0.15)
rrect(s, chart_rx, chart_y, chart_rw, chart_h, WHITE, LIGHT_GRAY, 0.5)
txt(s, "Kinerja Cabang", chart_rx+Inches(0.1), chart_y+Inches(0.08),
    chart_rw-Inches(0.2), Inches(0.3), sz=10, bold=True, color=DARK_TEXT)
branches = [("Sudirman", 94), ("Kelapa Gading", 87), ("BSD City", 82), ("Bekasi", 78)]
for j, (br, sc) in enumerate(branches):
    by3 = chart_y + Inches(0.5) + j * Inches(0.7)
    txt(s, br, chart_rx+Inches(0.12), by3, chart_rw*0.55, Inches(0.28), sz=9, color=DARK_TEXT)
    txt(s, f"{sc}%", chart_rx+chart_rw-Inches(0.45), by3, Inches(0.4), Inches(0.28),
        sz=9, bold=True, color=TEAL_DARK, align=PP_ALIGN.RIGHT)
    bw3 = (chart_rw-Inches(0.35)) * sc/100
    rect(s, chart_rx+Inches(0.12), by3+Inches(0.3), bw3, Pt(8), TEAL)
print("✓ Slide 15: Management Dashboard")

# ─────────────────────────────────────────
# SLIDE 16 — BUSINESS INTELLIGENCE
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Business Intelligence & Analytics Platform",
            "Transformasi data operasional menjadi insight strategis yang mendorong pertumbuhan bisnis")
footer(s)

bi_items = [
    ("📈", "Revenue & Trend Analytics", "MAGENTA",
     "Analisis tren revenue harian, mingguan, bulanan dengan drill-down per cabang dan per dokter. Deteksi anomali otomatis, perbandingan Year-over-Year, dan proyeksi 6 bulan berbasis machine learning."),
    ("👥", "Patient Behavior Analytics", "TEAL",
     "Analisis mendalam perilaku pasien: frekuensi kunjungan, jenis perawatan favorit, average transaction value, lifetime value, dan segmentasi untuk kampanye marketing yang hyper-targeted."),
    ("🦷", "Treatment & Service Analytics", "NAVY",
     "Performa tiap jenis layanan: volume, revenue, margin, dan satisfaction score. Identifikasi layanan underperforming. Benchmark antar cabang untuk best practice sharing yang berkelanjutan."),
    ("🏆", "Branch Performance Scorecard", "TEAL_DARK",
     "Scorecard komprehensif setiap cabang: NPS, revenue, efisiensi, utilisasi dokter, dan pertumbuhan pasien baru. Leaderboard untuk healthy competition dan identifikasi cabang yang butuh intervensi."),
    ("🔮", "AI-Powered Forecasting", "MAGENTA",
     "Prediksi demand pasien per cabang untuk optimasi resource dan staffing. Rekomendasi inventory supplies. Deteksi early churn signals untuk re-engagement proaktif sebelum pasien pergi ke kompetitor."),
    ("📊", "Real-Time Operational KPIs", "NAVY",
     "Dashboard KPI operasional 24/7: kapasitas antrian, utilisasi ruang perawatan, throughput pasien, dan revenue per jam. Alert otomatis jika KPI menyimpang dari target yang sudah ditetapkan."),
]

cw16 = (W - Inches(1.0)) / 3
ch16 = (H - ht - Inches(0.6)) / 2
for i, (ico, ttl, ac_name, dsc) in enumerate(bi_items):
    col = i % 3; row = i // 3
    fx16 = Inches(0.4) + col * (cw16 + Inches(0.1))
    fy16 = ht + Inches(0.15) + row * (ch16 + Inches(0.1))
    rrect(s, fx16, fy16, cw16, ch16, WHITE, LIGHT_GRAY, 0.5)
    ac16 = {"MAGENTA": MAGENTA, "TEAL": TEAL, "NAVY": NAVY, "TEAL_DARK": TEAL_DARK}[ac_name]
    rect(s, fx16, fy16, cw16, Pt(4), ac16)
    txt(s, ico, fx16+Inches(0.15), fy16+Inches(0.1), Inches(0.45), Inches(0.45), sz=18)
    txt(s, ttl, fx16+Inches(0.7), fy16+Inches(0.1), cw16-Inches(0.85), Inches(0.42),
        sz=12, bold=True, color=DARK_TEXT)
    txt(s, dsc, fx16+Inches(0.15), fy16+Inches(0.62), cw16-Inches(0.3), ch16-Inches(0.78),
        sz=9, color=GRAY)
print("✓ Slide 16: Business Intelligence")

# ─────────────────────────────────────────
# SLIDE 17 — PROJECT ROADMAP
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Roadmap Eksekusi Proyek: 8 Bulan Transformasi",
            "Pendekatan bertahap dan terstruktur yang meminimalkan risiko dan memaksimalkan kualitas delivery")
footer(s)

phases17 = [
    ("Bulan 1", "Discovery & Planning", MAGENTA, [
        "Kickoff & stakeholder alignment", "Audit sistem existing", "User research & persona", "Requirement final"]),
    ("Bulan 2", "UI/UX Design", TEAL_DARK, [
        "Design system & brand guide", "Wireframing semua platform", "Prototype & user testing", "Design approval"]),
    ("Bulan 3–4", "Website + SEO", NAVY, [
        "Development website", "Technical SEO setup", "Konten & copywriting", "Launch & monitoring"]),
    ("Bulan 4–6", "Mobile Development", RGBColor(0x7C,0x3A,0xED), [
        "PWA development", "React Native app", "Backend API & DB", "Payment & notifikasi"]),
    ("Bulan 6–7", "eKiosk Development", TEAL, [
        "Software kiosk", "Integrasi hardware", "Queue management", "Payment kiosk"]),
    ("Bulan 8", "Testing & Launch", RGBColor(0x10,0xB9,0x81), [
        "System integration test", "UAT bersama OMDC", "Staff training", "Go-Live! 🚀"]),
]

ph17 = (H - ht - Inches(0.65)) - Inches(0.1)
pw17 = (W - Inches(0.8)) / len(phases17)
gx = Inches(0.4)
gy = ht + Inches(0.15)

for i, (period, title, fc, tasks) in enumerate(phases17):
    px17 = gx + i * pw17
    # Month header
    rrect(s, px17+Pt(3), gy, pw17-Pt(6), Inches(0.6), fc)
    txt(s, period, px17+Pt(3), gy+Inches(0.05), pw17-Pt(6), Inches(0.28),
        sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, px17+Pt(3), gy+Inches(0.3), pw17-Pt(6), Inches(0.28),
        sz=9, color=RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER)
    # Task list
    task_y = gy + Inches(0.68)
    for j, task in enumerate(tasks):
        ty = task_y + j * Inches(0.6)
        rrect(s, px17+Pt(4), ty, pw17-Pt(8), Inches(0.52), WHITE, LIGHT_GRAY, 0.5)
        rect(s, px17+Pt(4), ty, Pt(4), Inches(0.52), fc)
        txt(s, task, px17+Inches(0.15), ty+Inches(0.08), pw17-Inches(0.3), Inches(0.38),
            sz=9, color=DARK_TEXT)
    # Bottom duration
    dy = task_y + 4 * Inches(0.6) + Inches(0.05)
    rect(s, px17+Pt(3), dy, pw17-Pt(6), Pt(4), fc)
print("✓ Slide 17: Project Roadmap")

# ─────────────────────────────────────────
# SLIDE 18 — EXPECTED IMPACT
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Dampak Bisnis yang Diproyeksikan",
            "KPI terukur yang akan dicapai dalam 12 bulan pertama pasca-implementasi penuh")
footer(s)

impacts = [
    ("40%", "Pengurangan\nWaktu Registrasi", "Dari 15 menit → 9 menit. eKiosk self-service menghilangkan bottleneck front desk sepenuhnya.", MAGENTA),
    ("35%", "Peningkatan\nKonversi Appointment", "Website SEO-optimized + booking frictionless mengubah lebih banyak pengunjung jadi pasien.", TEAL_DARK),
    ("50%", "Pengurangan\nAdmin Workload", "Otomasi registrasi, reminder, payment & EMR membebaskan staff untuk layanan bernilai tinggi.", NAVY),
    ("25%", "Peningkatan\nRepeat Visit", "Loyalty program digital, reminder perawatan berkala, dan engagement pasca-kunjungan yang personal.", RGBColor(0x7C,0x3A,0xED)),
]

cw18 = (W - Inches(1.0)) / 4
ch18 = H - ht - Inches(0.7)

for i, (num, lbl, desc, fc) in enumerate(impacts):
    ix = Inches(0.4) + i * (cw18 + Inches(0.06))
    iy = ht + Inches(0.2)
    rrect(s, ix, iy, cw18, ch18, WHITE, LIGHT_GRAY, 0.5)
    rect(s, ix, iy, cw18, Pt(5), fc)
    # Big number
    txt(s, num, ix, iy+Inches(0.2), cw18, Inches(1.1),
        sz=64, bold=True, color=fc, align=PP_ALIGN.CENTER)
    # Divider
    rect(s, ix+Inches(0.3), iy+Inches(1.35), cw18-Inches(0.6), Pt(2), LIGHT_GRAY)
    # Label
    txt(s, lbl, ix, iy+Inches(1.45), cw18, Inches(0.7),
        sz=14, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    # Description
    txt(s, desc, ix+Inches(0.15), iy+Inches(2.25), cw18-Inches(0.3), ch18-Inches(2.45),
        sz=10, color=GRAY, align=PP_ALIGN.CENTER)
    # Bottom badge
    rrect(s, ix+Inches(0.2), iy+ch18-Inches(0.5), cw18-Inches(0.4), Inches(0.38), fc)
    txt(s, "12 Bulan Pertama", ix+Inches(0.2), iy+ch18-Inches(0.5),
        cw18-Inches(0.4), Inches(0.38), sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
print("✓ Slide 18: Expected Impact")

# ─────────────────────────────────────────
# SLIDE 19 — INVESTMENT SUMMARY
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Rincian Investasi Software: Transparan & Kompetitif",
            "Setiap komponen dirancang memberikan nilai maksimal dengan investasi yang jelas dan terukur")
footer(s)

# Left: pricing table
lw19 = W * 0.55
items19 = [
    ("🌐", "Website Revamp Premium", "Multipage responsive, CMS, WhatsApp API, booking integration", "Rp 60.000.000", MAGENTA),
    ("🔍", "SEO & Digital Marketing", "Technical SEO, content strategy, Google Ads setup, 6-bln monitoring", "Rp 25.000.000", TEAL),
    ("📱", "Progressive Web App (PWA)", "Full-featured PWA dengan booking, membership, payment, notifikasi", "Rp 120.000.000", NAVY),
    ("🤖", "Android Native App", "React Native, App Store ready, offline mode, push notification", "Rp 90.000.000", RGBColor(0x7C,0x3A,0xED)),
    ("🍎", "iOS Native App", "React Native, TestFlight, App Store distribution, Apple Pay", "Rp 90.000.000", RGBColor(0x7C,0x3A,0xED)),
    ("🖥️", "eKiosk Software Platform", "Queue, payment, check-in, QR, real-time sync, monitoring dashboard", "Rp 75.000.000", TEAL_DARK),
]

row_h = (H - ht - Inches(1.1)) / (len(items19) + 1)
for i, (ico, ttl, dsc, price, fc) in enumerate(items19):
    ry = ht + Inches(0.15) + i * row_h
    rrect(s, Inches(0.4), ry+Pt(2), lw19-Inches(0.2), row_h-Pt(4), WHITE, LIGHT_GRAY, 0.5)
    rect(s, Inches(0.4), ry+Pt(2), Pt(4), row_h-Pt(4), fc)
    txt(s, ico, Inches(0.55), ry+Inches(0.05), Inches(0.4), row_h-Pt(8), sz=16, align=PP_ALIGN.CENTER)
    txt(s, ttl, Inches(1.05), ry+Inches(0.05), lw19-Inches(2.8), Inches(0.32),
        sz=12, bold=True, color=DARK_TEXT)
    txt(s, dsc, Inches(1.05), ry+Inches(0.35), lw19-Inches(2.8), row_h-Inches(0.5),
        sz=9, color=GRAY)
    txt(s, price, lw19-Inches(1.55), ry+Inches(0.1), Inches(1.3), row_h-Pt(8),
        sz=13, bold=True, color=fc, align=PP_ALIGN.RIGHT)

# Subtotal row
st_y = ht + Inches(0.15) + len(items19) * row_h
rrect(s, Inches(0.4), st_y+Pt(3), lw19-Inches(0.2), row_h-Pt(6), NAVY)
txt(s, "TOTAL SOFTWARE", Inches(0.55), st_y+Inches(0.07), lw19-Inches(2.5), row_h-Inches(0.2),
    sz=13, bold=True, color=WHITE)
txt(s, "Rp 460.000.000", lw19-Inches(1.55), st_y+Inches(0.07), Inches(1.3), row_h-Inches(0.2),
    sz=15, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)

# Right: value summary cards
rx19 = lw19 + Inches(0.2)
rw19 = W - rx19 - Inches(0.4)
vals19 = [
    ("6", "Modul Platform\nLengkap", MAGENTA),
    ("25", "Slides Deliverable\nProposal Ini", TEAL),
    ("8", "Bulan Timeline\nPengerjaan", NAVY),
    ("3", "Bulan Support\nPost-Launch", TEAL_DARK),
]
vch = (H - ht - Inches(0.6)) / 4
for i, (num, lbl, fc) in enumerate(vals19):
    vy = ht + Inches(0.15) + i * vch
    rrect(s, rx19, vy+Pt(3), rw19, vch-Pt(6), WHITE, LIGHT_GRAY, 0.5)
    rect(s, rx19, vy+Pt(3), Pt(4), vch-Pt(6), fc)
    txt(s, num, rx19, vy+Inches(0.05), rw19, vch*0.55,
        sz=32, bold=True, color=fc, align=PP_ALIGN.CENTER)
    txt(s, lbl, rx19, vy+vch*0.55, rw19, vch*0.38,
        sz=10, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
print("✓ Slide 19: Investment Summary")

# ─────────────────────────────────────────
# SLIDE 20 — HARDWARE INVESTMENT
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Investasi Hardware & Total Nilai Proyek",
            "3 unit Smart eKiosk enterprise-grade untuk transformasi fisik pengalaman pasien di klinik")
footer(s)

# 3 Kiosk units
unit_w = Inches(2.8)
unit_h = Inches(3.5)
unit_y = ht + Inches(0.2)
for i in range(3):
    ux = Inches(0.5) + i * Inches(3.0)
    rrect(s, ux, unit_y, unit_w, unit_h, NAVY, MAGENTA, 1)
    txt(s, "🖥️", ux, unit_y+Inches(0.15), unit_w, Inches(0.8),
        sz=40, align=PP_ALIGN.CENTER, color=WHITE)
    txt(s, f"UNIT {i+1:02d}", ux, unit_y+Inches(1.0), unit_w, Inches(0.38),
        sz=12, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)
    txt(s, f"Cabang {['Sudirman','Kelapa Gading','BSD City'][i]}",
        ux, unit_y+Inches(1.38), unit_w, Inches(0.32),
        sz=10, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Rp 28.000.000", ux, unit_y+Inches(1.78), unit_w, Inches(0.45),
        sz=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    specs_short = ["32\" Touchscreen", "Intel i5 Industrial", "Thermal Printer", "QR Scanner"]
    for j, sp in enumerate(specs_short):
        ty2 = unit_y + Inches(2.3) + j * Inches(0.28)
        txt(s, f"✓  {sp}", ux+Inches(0.25), ty2, unit_w-Inches(0.3), Inches(0.25),
            sz=8, color=RGBColor(0xB0,0xCC,0xE8))

    # Label
    txt(s, f"Termasuk: Instalasi & Konfigurasi",
        ux, unit_y+unit_h+Inches(0.08), unit_w, Inches(0.28),
        sz=8, color=GRAY, align=PP_ALIGN.CENTER)

# Connector + arrows
for i in range(2):
    ax20 = Inches(0.5) + (i+1) * Inches(3.0) - Inches(0.2)
    txt(s, "+", ax20, unit_y + unit_h/2, Inches(0.4), Inches(0.4),
        sz=20, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)

# Summary box
sum_y = unit_y + unit_h + Inches(0.5)
sum_h = H - sum_y - Inches(0.5)
col_s = (W - Inches(0.8)) / 3
totals = [
    ("Software Total", "Rp 460.000.000", NAVY),
    ("Hardware 3 Unit", "Rp 84.000.000", TEAL_DARK),
    ("TOTAL INVESTASI", "Rp 544.000.000", MAGENTA),
]
for i, (lbl, val, fc) in enumerate(totals):
    tx2 = Inches(0.4) + i * col_s
    rrect(s, tx2+Pt(4), sum_y, col_s-Pt(8), sum_h, fc if i==2 else WHITE,
          LIGHT_GRAY if i<2 else None, 0.5 if i<2 else None)
    if i == 2: rect(s, tx2+Pt(4), sum_y, col_s-Pt(8), Pt(5), WHITE)
    txt(s, lbl, tx2+Pt(4), sum_y+Inches(0.08), col_s-Pt(8), Inches(0.35),
        sz=11, bold=True, color=WHITE if i==2 else DARK_TEXT, align=PP_ALIGN.CENTER)
    txt(s, val, tx2+Pt(4), sum_y+Inches(0.42), col_s-Pt(8), sum_h-Inches(0.5),
        sz=16, bold=True, color=WHITE if i==2 else fc, align=PP_ALIGN.CENTER)
print("✓ Slide 20: Hardware Investment")

# ─────────────────────────────────────────
# SLIDE 21 — SPECIAL PILOT PROGRAM
# ─────────────────────────────────────────
s = blank_slide()
bg(s, NAVY)
oval(s, W-Inches(6), Inches(-2), Inches(8), Inches(8), NAVY2)
oval(s, Inches(-2), H-Inches(4), Inches(6), Inches(6), RGBColor(0x00,0x50,0x55))
rect(s, 0, 0, Inches(0.45), H, MAGENTA)
rect(s, 0, H-Inches(0.08), W, Inches(0.08), TEAL)

txt(s, "PENAWARAN EKSKLUSIF OMDC", Inches(0.8), Inches(0.5), W-Inches(1.2), Inches(0.5),
    sz=12, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)

# Crossed out price
rrect(s, Inches(2.5), Inches(1.0), W-Inches(5.0), Inches(1.0), NAVY2, GRAY, 0.5)
txt(s, "Rp 544.000.000", Inches(2.5), Inches(1.05), W-Inches(5.0), Inches(0.85),
    sz=28, bold=True, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
# Strikethrough line
rect(s, Inches(2.8), Inches(1.47), W-Inches(5.6), Pt(4), MAGENTA)

# Big arrow down
txt(s, "▼", Inches(0.8), Inches(2.1), W-Inches(1.6), Inches(0.7),
    sz=32, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

txt(s, "✦  SPECIAL PILOT PROGRAM  ✦", Inches(0.8), Inches(2.85), W-Inches(1.6), Inches(0.55),
    sz=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Main price
rrect(s, Inches(1.5), Inches(3.4), W-Inches(3.0), Inches(1.6), MAGENTA)
txt(s, "Rp 499.000.000", Inches(1.5), Inches(3.45), W-Inches(3.0), Inches(1.5),
    sz=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "HEMAT Rp 45 JUTA  •  All-Inclusive: Software + Hardware + Instalasi + Training + 3 Bln Support",
    Inches(0.8), Inches(5.1), W-Inches(1.6), Inches(0.42),
    sz=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

txt(s, "⏰  Penawaran berlaku hingga 31 Juli 2026 — Hanya untuk mitra pilot program Ruang.dev 2026",
    Inches(0.8), Inches(5.58), W-Inches(1.6), Inches(0.38),
    sz=11, color=GRAY, align=PP_ALIGN.CENTER)

txt(s, "Ruang.dev  •  hello@ruang.dev", Inches(0.8), Inches(6.1), W-Inches(1.6), Inches(0.38),
    sz=10, color=GRAY, align=PP_ALIGN.CENTER)
print("✓ Slide 21: Special Pilot Program")

# ─────────────────────────────────────────
# SLIDE 22 — PAYMENT MILESTONE
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Milestone Pembayaran Bertahap",
            "Struktur pembayaran adil dan transparan — selaras dengan progress pengerjaan yang nyata")
footer(s)

milestones = [
    ("30%", "Penandatanganan Kontrak", "Rp 149.700.000", "Kickoff resmi, pembentukan tim proyek, dan mulai fase discovery & planning.", MAGENTA),
    ("20%", "Persetujuan Desain", "Rp 99.800.000", "Desain UI/UX seluruh platform telah disetujui stakeholder. Milestone kreatif tercapai.", TEAL),
    ("20%", "Progress Development 50%", "Rp 99.800.000", "Development backend dan frontend mencapai 50%. Demo internal fitur-fitur utama.", NAVY),
    ("20%", "User Acceptance Testing", "Rp 99.800.000", "UAT selesai, semua bug kritis terselesaikan. Platform siap untuk go-live.", TEAL_DARK),
    ("10%", "Go-Live & Serah Terima", "Rp 49.900.000", "Launch resmi semua platform. Pelatihan staff. Handover dokumentasi teknis lengkap.", RGBColor(0x10,0xB9,0x81)),
]

mw = (W - Inches(0.8)) / len(milestones)
my_base = ht + Inches(0.2)
mh = H - ht - Inches(0.7)

# Timeline line
rect(s, Inches(0.4), my_base + mh/2, W-Inches(0.8), Pt(3), LIGHT_GRAY)

for i, (pct, phase, amount, desc, fc) in enumerate(milestones):
    mx22 = Inches(0.4) + i * mw
    # Alternating up/down cards
    if i % 2 == 0:
        cy22 = my_base
        ch22 = mh/2 - Inches(0.4)
    else:
        cy22 = my_base + mh/2 + Inches(0.35)
        ch22 = mh/2 - Inches(0.35)

    rrect(s, mx22+Pt(4), cy22, mw-Pt(8), ch22, WHITE, LIGHT_GRAY, 0.5)
    rect(s, mx22+Pt(4), cy22, mw-Pt(8), Pt(4), fc)
    txt(s, pct, mx22+Pt(4), cy22+Inches(0.05), mw-Pt(8), Inches(0.5),
        sz=26, bold=True, color=fc, align=PP_ALIGN.CENTER)
    txt(s, phase, mx22+Pt(8), cy22+Inches(0.52), mw-Pt(16), Inches(0.38),
        sz=10, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    txt(s, amount, mx22+Pt(4), cy22+Inches(0.88), mw-Pt(8), Inches(0.32),
        sz=11, bold=True, color=fc, align=PP_ALIGN.CENTER)
    txt(s, desc, mx22+Pt(8), cy22+Inches(1.22), mw-Pt(16), ch22-Inches(1.35),
        sz=8, color=GRAY, align=PP_ALIGN.CENTER)

    # Connector dot on timeline
    oval(s, mx22+mw/2-Inches(0.15), my_base+mh/2-Inches(0.15),
         Inches(0.3), Inches(0.3), fc)
    txt(s, str(i+1), mx22+mw/2-Inches(0.15), my_base+mh/2-Inches(0.15),
        Inches(0.3), Inches(0.3), sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
print("✓ Slide 22: Payment Milestone")

# ─────────────────────────────────────────
# SLIDE 23 — SUPPORT & MAINTENANCE
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Paket Dukungan & Pemeliharaan Berkelanjutan",
            "Kemitraan jangka panjang untuk memastikan sistem OMDC selalu prima, aman, dan berkembang")
footer(s)

# Header price card on top right
rrect(s, W-Inches(3.2), ht+Inches(0.1), Inches(2.8), Inches(1.1), MAGENTA)
txt(s, "Rp 10.000.000", W-Inches(3.2), ht+Inches(0.12), Inches(2.8), Inches(0.55),
    sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "per bulan  •  All-Inclusive", W-Inches(3.2), ht+Inches(0.65), Inches(2.8), Inches(0.38),
    sz=10, color=RGBColor(0xFC,0xE7,0xF3), align=PP_ALIGN.CENTER)

supports = [
    ("☁️", "Cloud Monitoring 24/7", MAGENTA, "Pemantauan infrastruktur dan aplikasi non-stop. Alert otomatis untuk anomali performa, downtime, atau ancaman keamanan. Response time insiden < 15 menit untuk isu critical."),
    ("🎧", "Technical Support", TEAL, "Tim support dedikasi via tiket, WhatsApp, dan telepon pada jam kerja. Eskalasi 24/7 untuk isu critical. SLA: response 2 jam, resolusi 24 jam untuk masalah level high."),
    ("🐛", "Bug Fix Priority", NAVY, "Penanganan bug prioritas dengan severity levels yang jelas. Hotfix deployment kapanpun dibutuhkan. Root cause analysis dan laporan insiden untuk setiap isu mayor."),
    ("🔒", "Security Updates", TEAL_DARK, "Patch keamanan rutin untuk semua komponen: OS, framework, dependencies, SSL/TLS renewal. Vulnerability scanning bulanan. Penetration testing tahunan."),
    ("🚀", "Minor Enhancement", RGBColor(0x7C,0x3A,0xED), "Quarterly release dengan minor feature enhancement berdasarkan feedback operasional. Akses priority roadmap. Sesi strategy review per kuartal bersama tim Ruang.dev."),
    ("📊", "Monthly Reporting", RGBColor(0x10,0xB9,0x81), "Laporan performa bulanan: uptime, response time, error rate, keamanan, dan rekomendasi optimasi. Business review deck siap untuk rapat direksi."),
]

cw23 = (W - Inches(1.0)) / 3
ch23 = (H - ht - Inches(1.4)) / 2

for i, (ico, ttl, fc, dsc) in enumerate(supports):
    col = i % 3; row = i // 3
    fx23 = Inches(0.4) + col * (cw23 + Inches(0.1))
    fy23 = ht + Inches(1.3) + row * (ch23 + Inches(0.1))
    rrect(s, fx23, fy23, cw23, ch23, WHITE, LIGHT_GRAY, 0.5)
    rect(s, fx23, fy23, cw23, Pt(4), fc)
    txt(s, ico, fx23+Inches(0.15), fy23+Inches(0.08), Inches(0.45), Inches(0.45), sz=18)
    txt(s, ttl, fx23+Inches(0.7), fy23+Inches(0.1), cw23-Inches(0.85), Inches(0.38),
        sz=12, bold=True, color=DARK_TEXT)
    txt(s, dsc, fx23+Inches(0.15), fy23+Inches(0.6), cw23-Inches(0.3), ch23-Inches(0.75),
        sz=9, color=GRAY)
print("✓ Slide 23: Support & Maintenance")

# ─────────────────────────────────────────
# SLIDE 24 — WHY RUANG.DEV
# ─────────────────────────────────────────
s = blank_slide()
bg(s)
ht = header(s, "Mengapa Ruang.dev: Partner Transformasi Digital Terpercaya",
            "Bukan sekadar vendor teknologi — kami adalah mitra strategis yang invested dalam kesuksesan OMDC")
footer(s)

pillars = [
    ("🎯", "End-to-End Delivery", MAGENTA,
     "Satu tim, satu tanggung jawab penuh dari discovery, design, development, testing, hingga launch dan maintenance. Tidak ada fragmentasi vendor—akuntabilitas penuh di satu tangan."),
    ("💡", "Product Thinking", TEAL,
     "Kami proaktif menyarankan fitur dan pendekatan yang lebih optimal. Mindset product manager, bukan sekadar order-taker. Kami bertanya 'mengapa' sebelum langsung mengerjakan 'apa'."),
    ("🏥", "Healthcare Experience", NAVY,
     "Pengalaman nyata membangun sistem digital untuk institusi kesehatan Indonesia. Memahami regulasi, alur klinis, kebutuhan privasi data pasien, dan tantangan unik digitalisasi layanan medis."),
    ("⚙️", "Scalable Architecture", TEAL_DARK,
     "Setiap sistem yang kami bangun menggunakan cloud-native, microservices architecture yang siap scale dari 1 cabang ke 100 cabang tanpa major rearchitecture. Future-proof dari hari pertama."),
    ("🤝", "Long-Term Partnership", RGBColor(0x7C,0x3A,0xED),
     "Visi kami adalah tumbuh bersama OMDC selama 5–10 tahun ke depan. Quarterly strategy review, product roadmap bersama, priority access ke fitur baru, dan tarif preferensial untuk proyek berikutnya."),
    ("✅", "Quality & Delivery", RGBColor(0x10,0xB9,0x81),
     "Proses development dengan automated testing, code review ketat, CI/CD pipeline, dan quality gate di setiap milestone. Tidak ada delivery tanpa UAT sign-off dari stakeholder OMDC."),
]

cw24 = (W - Inches(1.0)) / 3
ch24 = (H - ht - Inches(0.6)) / 2

for i, (ico, ttl, fc, dsc) in enumerate(pillars):
    col = i % 3; row = i // 3
    fx24 = Inches(0.4) + col * (cw24 + Inches(0.1))
    fy24 = ht + Inches(0.15) + row * (ch24 + Inches(0.1))
    # Card with accent top
    rrect(s, fx24, fy24, cw24, ch24, WHITE, LIGHT_GRAY, 0.5)
    rect(s, fx24, fy24, cw24, Pt(5), fc)
    # Icon badge
    rrect(s, fx24+Inches(0.15), fy24+Inches(0.12), Inches(0.6), Inches(0.6),
          RGBColor(0xFC,0xE7,0xF3))
    txt(s, ico, fx24+Inches(0.15), fy24+Inches(0.12), Inches(0.6), Inches(0.6),
        sz=22, align=PP_ALIGN.CENTER)
    txt(s, ttl, fx24+Inches(0.9), fy24+Inches(0.14), cw24-Inches(1.05), Inches(0.42),
        sz=13, bold=True, color=DARK_TEXT)
    txt(s, dsc, fx24+Inches(0.15), fy24+Inches(0.82), cw24-Inches(0.3), ch24-Inches(0.98),
        sz=10, color=GRAY)
print("✓ Slide 24: Why Ruang.dev")

# ─────────────────────────────────────────
# SLIDE 25 — THANK YOU
# ─────────────────────────────────────────
s = blank_slide()
bg(s, NAVY)
oval(s, W-Inches(5), Inches(-1), Inches(6), Inches(6), NAVY2)
oval(s, Inches(-1.5), H-Inches(3), Inches(5), Inches(5), RGBColor(0x00,0x50,0x55))
rect(s, 0, 0, Inches(0.45), H, MAGENTA)
rect(s, 0, H-Inches(0.08), W, Inches(0.08), TEAL)

txt(s, "OMDC × Ruang.dev", Inches(0.8), Inches(0.6), W-Inches(1.6), Inches(0.5),
    sz=14, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)

txt(s, "Membangun Masa Depan", Inches(0.8), Inches(1.3), W-Inches(1.6), Inches(0.9),
    sz=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Digital Healthcare Indonesia", Inches(0.8), Inches(2.1), W-Inches(1.6), Inches(0.85),
    sz=46, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)

rect(s, Inches(3.5), Inches(3.1), W-Inches(7.0), Pt(2), RGBColor(0x2D,0x3F,0x7A))

txt(s, "Bersama, kita wujudkan klinik gigi paling digital dan patient-centric di Indonesia",
    Inches(0.8), Inches(3.3), W-Inches(1.6), Inches(0.5),
    sz=14, color=RGBColor(0xA0,0xB4,0xD0), align=PP_ALIGN.CENTER)

# Contact info cards
contacts = [
    ("🌐", "ruang.dev"),
    ("📧", "hello@ruang.dev"),
    ("📱", "+62 812-0000-0000"),
    ("📍", "Jakarta, Indonesia"),
]
total_cw = Inches(2.4) * len(contacts)
cstart = (W - total_cw) / 2
for i, (ico, ctxt) in enumerate(contacts):
    cx25 = cstart + i * Inches(2.6)
    rrect(s, cx25, Inches(4.05), Inches(2.4), Inches(0.62), NAVY2, TEAL, 0.5)
    txt(s, f"{ico}  {ctxt}", cx25, Inches(4.1), Inches(2.4), Inches(0.5),
        sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Siap memulai transformasi? Mari diskusi lebih lanjut.",
    Inches(0.8), Inches(4.85), W-Inches(1.6), Inches(0.4),
    sz=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Ecosystem reminder
eco_items = ["Website Premium", "PWA + Mobile App", "Smart eKiosk", "Analytics Dashboard"]
ew = (W - Inches(1.0)) / len(eco_items)
ey = Inches(5.4)
for i, item in enumerate(eco_items):
    ex25 = Inches(0.4) + i * ew
    rrect(s, ex25+Pt(4), ey, ew-Pt(8), Inches(0.55), NAVY2, MAGENTA, 0.5)
    txt(s, item, ex25+Pt(4), ey+Inches(0.08), ew-Pt(8), Inches(0.38),
        sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Rp 499.000.000 — Special Pilot Program 2026",
    Inches(0.8), Inches(6.15), W-Inches(1.6), Inches(0.42),
    sz=12, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)
print("✓ Slide 25: Thank You")

# ─────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────
output = "/home/user/2kang/OMDC_Digital_Transformation_2026.pptx"
prs.save(output)
print(f"\n✅ Presentation saved: {output}")
print(f"   Total slides: {len(prs.slides)}")
